import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image

st.set_page_config(page_title="PowerPoint Merger by OrvilleDev", layout="centered")

# Custom CSS to style file uploader buttons to red
st.markdown("""
<style>
    /* Style file uploader "Browse files" buttons to be red */
    div[data-testid="stFileUploader"] button {
        background-color: #FF0000 !important;
        color: white !important;
        border: none !important;
        border-radius: 0.25rem !important;
    }
    div[data-testid="stFileUploader"] button:hover {
        background-color: #CC0000 !important;
    }
    div[data-testid="stFileUploader"] button:focus {
        background-color: #CC0000 !important;
    }
    div[data-testid="stFileUploader"] button:active {
        background-color: #AA0000 !important;
    }
</style>
""", unsafe_allow_html=True)

st.title("📊 PowerPoint Merger")

def resize_image_to_1920x1080(image_bytes):
    """Resize image to 1920x1080 pixels"""
    try:
        # Open image from bytes
        img = Image.open(BytesIO(image_bytes))
        # Resize to 1920x1080 (use LANCZOS resampling for quality)
        try:
            # Try newer API first
            resized_img = img.resize((1920, 1080), Image.Resampling.LANCZOS)
        except AttributeError:
            # Fallback for older Pillow versions
            resized_img = img.resize((1920, 1080), Image.LANCZOS)
        # Convert to RGB if necessary (for JPEG compatibility)
        if resized_img.mode != 'RGB':
            resized_img = resized_img.convert('RGB')
        # Save to bytes
        output = BytesIO()
        resized_img.save(output, format='PNG')
        output.seek(0)
        return output.getvalue()
    except Exception as e:
        # If resizing fails, return original
        return image_bytes

def extract_text_from_slide(slide):
    """Extract all text from a slide"""
    text_content = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = ""
                for run in paragraph.runs:
                    para_text += run.text
                if para_text.strip():
                    text_content.append(para_text.strip())
    return "\n".join(text_content)

def is_all_caps(text):
    """Check if text is all uppercase (all caps)"""
    return text.isupper() and any(c.isalpha() for c in text)

def parse_txt_file(txt_content):
    """Parse .txt file and extract slides with titles and verses"""
    lines = txt_content.decode('utf-8').split('\n')
    slides = []
    current_slide = {'title': None, 'verses': []}
    
    for line in lines:
        line = line.strip()
        
        # Check if line is a title
        if line.upper().startswith('TITLE:'):
            # If we have content in current slide, save it
            if current_slide['title'] or current_slide['verses']:
                slides.append(current_slide)
            
            # Extract title text (handle both TITLE:text and TITLE:"text" formats)
            title_text = line[6:].strip()  # Remove "TITLE:"
            if title_text.startswith('"') and title_text.endswith('"'):
                title_text = title_text[1:-1]  # Remove quotes
            
            # Start new slide with this title
            current_slide = {'title': title_text, 'verses': []}
        elif line == '':
            # Blank line - if we have content, save current slide and start new one
            if current_slide['title'] or current_slide['verses']:
                slides.append(current_slide)
                current_slide = {'title': None, 'verses': []}
        else:
            # Regular verse line
            if line:  # Only add non-empty lines
                current_slide['verses'].append(line)
    
    # Add the last slide if it has content
    if current_slide['title'] or current_slide['verses']:
        slides.append(current_slide)
    
    return slides

def create_formatted_slide(target_presentation, text, is_title, title_color, verse_color, title_font_size, verse_font_size, title_font, verse_font, background_image=None):
    """Create a new slide with formatted text"""
    # Use blank layout
    blank_slide_layout = target_presentation.slide_layouts[6]
    slide = target_presentation.slides.add_slide(blank_slide_layout)
    
    # Get slide dimensions
    slide_width = target_presentation.slide_width
    slide_height = target_presentation.slide_height
    
    # Add background image if provided (add it first so text appears on top)
    if background_image:
        try:
            # Add image to cover entire slide
            slide.shapes.add_picture(
                BytesIO(background_image),
                0,  # left
                0,  # top
                slide_width,  # width
                slide_height  # height
            )
        except Exception:
            pass
    
    # Set black background (if no image, or as fallback)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
    
    # Create text box stretching from left to right edge
    width = slide_width  # Full width of slide
    height = Inches(7)  # Keep reasonable height
    left = 0  # Start from leftmost edge
    top = (slide_height - height) / 2  # Center vertically
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Middle vertical alignment
    text_frame.margin_left = Inches(0.5)  # Small padding from left edge
    text_frame.margin_right = Inches(0.5)  # Small padding from right edge
    text_frame.margin_top = Inches(0.5)
    text_frame.margin_bottom = Inches(0.5)
    
    # Clear default paragraph
    text_frame.clear()
    
    # Add text with formatting
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER  # Center alignment
    paragraph.space_after = Pt(0)
    
    run = paragraph.add_run()
    run.text = text  # Preserve original case
    font = run.font
    
    # Check if text is all uppercase (all caps)
    is_all_caps = text.isupper() and any(c.isalpha() for c in text)
    
    # Set font properties: use title settings if it's marked as title OR if it's all caps
    if is_title or is_all_caps:
        font.name = title_font  # Use selected title font
        font.size = Pt(title_font_size)  # Use selected title font size
        font.color.rgb = RGBColor(*title_color)  # Use selected title color
        font.bold = True  # Bold for titles
    else:
        font.name = verse_font  # Use selected verse font
        font.size = Pt(verse_font_size)  # Use selected verse font size
        font.color.rgb = RGBColor(*verse_color)  # Use selected verse color
        font.bold = False
    
    return slide

def create_template_powerpoint(title_color, verse_color, title_font_size, verse_font_size, title_font, verse_font, background_image=None):
    """Create a PowerPoint template with 1 title slide and 1 verse slide"""
    template_prs = Presentation()
    
    # Set slide dimensions to 16:9 Widescreen aspect ratio
    template_prs.slide_width = Inches(13.33)
    template_prs.slide_height = Inches(7.5)
    
    # Remove default empty slide
    if template_prs.slides:
        template_prs.slides.remove(template_prs.slides[0])
    
    # Create title slide
    create_formatted_slide(template_prs, "YOUR TITLE HERE", True, title_color, verse_color, title_font_size, verse_font_size, title_font, verse_font, background_image)
    
    # Create verse slide
    create_formatted_slide(template_prs, "Your verse text here\nYou can add multiple lines\nEach line will appear on the slide", False, title_color, verse_color, title_font_size, verse_font_size, title_font, verse_font, background_image)
    
    return template_prs

def create_template_txt():
    """Create a .txt template with TITLE: format"""
    template_content = """TITLE: Your Title Here

Your verse text here
You can add multiple lines
Each line will appear on the slide

TITLE: Another Title (Optional)

More verse text here
Add as many titles and verses as you need
Separate slides with blank lines"""
    return template_content

# Initialize session state for file ordering
if 'file_order' not in st.session_state:
    st.session_state.file_order = []
if 'uploaded_files_dict' not in st.session_state:
    st.session_state.uploaded_files_dict = {}
if 'title_color' not in st.session_state:
    st.session_state.title_color = [255, 255, 0]  # Default yellow
if 'verse_color' not in st.session_state:
    st.session_state.verse_color = [255, 255, 255]  # Default white
if 'title_font_size' not in st.session_state:
    st.session_state.title_font_size = 72  # Default 72pt
if 'verse_font_size' not in st.session_state:
    st.session_state.verse_font_size = 65  # Default 65pt
if 'title_font' not in st.session_state:
    st.session_state.title_font = 'Arial'  # Default Arial
if 'verse_font' not in st.session_state:
    st.session_state.verse_font = 'Arial'  # Default Arial
if 'background_image' not in st.session_state:
    st.session_state.background_image = None
if 'txt_files_dict' not in st.session_state:
    st.session_state.txt_files_dict = {}

# Common fonts compatible across all systems
COMMON_FONTS = [
    'Arial',
    'Times New Roman',
    'Calibri',
    'Verdana',
    'Georgia',
    'Tahoma',
    'Trebuchet MS',
    'Courier New',
    'Comic Sans MS',
    'Impact'
]

# Color and Font Settings
st.subheader("Color Settings")
col1, col2 = st.columns(2)

with col1:
    # Convert RGB to hex for color picker
    title_hex = f"#{st.session_state.title_color[0]:02x}{st.session_state.title_color[1]:02x}{st.session_state.title_color[2]:02x}".upper()
    title_color = st.color_picker("Title Color (All Caps)", title_hex, key="title_color_picker")
    # Convert hex to RGB
    title_color_rgb = tuple(int(title_color[i:i+2], 16) for i in (1, 3, 5))
    st.session_state.title_color = list(title_color_rgb)

with col2:
    # Convert RGB to hex for color picker
    verse_hex = f"#{st.session_state.verse_color[0]:02x}{st.session_state.verse_color[1]:02x}{st.session_state.verse_color[2]:02x}".upper()
    verse_color = st.color_picker("Verse Color", verse_hex, key="verse_color_picker")
    # Convert hex to RGB
    verse_color_rgb = tuple(int(verse_color[i:i+2], 16) for i in (1, 3, 5))
    st.session_state.verse_color = list(verse_color_rgb)

st.subheader("Font Size Settings")
font_col1, font_col2 = st.columns(2)

with font_col1:
    title_font_size = st.number_input("Title Font Size (pt)", min_value=10, max_value=200, value=st.session_state.title_font_size, step=1, key="title_font_size_input")
    st.session_state.title_font_size = int(title_font_size)

with font_col2:
    verse_font_size = st.number_input("Verse Font Size (pt)", min_value=10, max_value=200, value=st.session_state.verse_font_size, step=1, key="verse_font_size_input")
    st.session_state.verse_font_size = int(verse_font_size)

st.subheader("Font Family Settings")
font_family_col1, font_family_col2 = st.columns(2)

with font_family_col1:
    title_font = st.selectbox("Title Font", COMMON_FONTS, index=COMMON_FONTS.index(st.session_state.title_font) if st.session_state.title_font in COMMON_FONTS else 0, key="title_font_select")
    st.session_state.title_font = title_font

with font_family_col2:
    verse_font = st.selectbox("Verse Font", COMMON_FONTS, index=COMMON_FONTS.index(st.session_state.verse_font) if st.session_state.verse_font in COMMON_FONTS else 0, key="verse_font_select")
    st.session_state.verse_font = verse_font

st.markdown("## Download Templates")
st.write("Download templates to get started. Edit them and upload back to the app.")
col_template_pptx, col_template_txt = st.columns(2)

with col_template_pptx:
    # Generate template data
    try:
        template_prs = create_template_powerpoint(
            st.session_state.title_color,
            st.session_state.verse_color,
            st.session_state.title_font_size,
            st.session_state.verse_font_size,
            st.session_state.title_font,
            st.session_state.verse_font,
            st.session_state.background_image
        )
        template_output = BytesIO()
        template_prs.save(template_output)
        template_output.seek(0)
        pptx_template_data = template_output.getvalue()
        
        st.download_button(
            label="📥 Download PowerPoint Template",
            data=pptx_template_data,
            file_name="powerpoint_template.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="download_pptx_template_file",
            type="primary"
        )
    except Exception as e:
        st.error(f"Error creating PowerPoint template: {str(e)}")

with col_template_txt:
    # Generate template data
    try:
        template_txt = create_template_txt()
        txt_template_data = template_txt.encode('utf-8')
        
        st.download_button(
            label="📥 Download Text Template",
            data=txt_template_data,
            file_name="text_template.txt",
            mime="text/plain",
            key="download_txt_template_file",
            type="primary"
        )
    except Exception as e:
        st.error(f"Error creating text template: {str(e)}")

st.divider()

st.markdown("## Upload PowerPoint Files")
uploaded_files = st.file_uploader(
    "Upload PowerPoint files",
    type=["pptx"],
    accept_multiple_files=True,
    key="pptx_file_uploader"
)

st.markdown("## Upload Text Files (Optional)")
uploaded_txt_files = st.file_uploader(
    "Upload .txt files with TITLE: format",
    type=["txt"],
    accept_multiple_files=True,
    key="txt_file_uploader"
)

st.subheader("Background Image (Optional)")
background_image_file = st.file_uploader(
    "Upload background image (will be applied to all slides)",
    type=["png", "jpg", "jpeg", "gif", "bmp"],
    key="background_image_uploader"
)

if background_image_file:
    # Read image data
    image_data = background_image_file.read()
    # Resize to 1920x1080
    resized_image_data = resize_image_to_1920x1080(image_data)
    st.session_state.background_image = resized_image_data
    st.success(f"✅ Background image loaded and resized to 1920x1080: {background_image_file.name}")
    # Display image preview
    background_image_file.seek(0)  # Reset file pointer
    st.image(background_image_file, width=300)
    if st.button("Remove Background Image", key="remove_bg_image"):
        st.session_state.background_image = None
        st.rerun()
else:
    st.session_state.background_image = None

# Update session state when new files are uploaded
# Store file bytes in session state so files persist across reruns
if uploaded_files is not None and len(uploaded_files) > 0:
    # Process newly uploaded files
    for file in uploaded_files:
        # Read ALL file bytes and store in session state
        # Make sure we read from the beginning
        file.seek(0)
        file_bytes = file.read()
        # Create a copy of the bytes to ensure we have all the data
        file_bytes = bytes(file_bytes)
        
        # Store file info: name and bytes
        file_info = {
            'name': file.name,
            'bytes': file_bytes,
            'type': 'pptx'
        }
        st.session_state.uploaded_files_dict[file.name] = file_info
        
        # Add to order if not already present
        if file.name not in st.session_state.file_order:
            st.session_state.file_order.append(file.name)
    
    # Don't automatically remove files - let the user remove them via the Remove button
    # Files persist in session state even if not in current upload
    # This allows multiple files to be uploaded and kept

# Also handle txt files in the ordering
# Store file bytes in session state so files persist across reruns
if uploaded_txt_files is not None and len(uploaded_txt_files) > 0:
    # Process newly uploaded txt files
    for file in uploaded_txt_files:
        # Read file bytes and store in session state
        file.seek(0)  # Reset file pointer
        file_bytes = file.read()
        # Create a copy of the bytes to ensure we have all the data
        file_bytes = bytes(file_bytes)
        
        # Store file info: name and bytes
        file_info = {
            'name': file.name,
            'bytes': file_bytes,
            'type': 'txt'
        }
        st.session_state.txt_files_dict[file.name] = file_info
        
        # Add to order if not already present
        if file.name not in st.session_state.file_order:
            st.session_state.file_order.append(file.name)
    
    # Don't automatically remove files - let the user remove them via the Remove button
    # Files persist in session state even if not in current upload

# Display file reordering interface
if st.session_state.file_order:
    st.subheader("Arrange Files (use buttons to reorder)")
    
    # Filter to only show valid files (remove any manual slide references)
    valid_file_order = [item_id for item_id in st.session_state.file_order 
                       if not item_id.startswith("MANUAL_SLIDE_") and 
                       (item_id in st.session_state.uploaded_files_dict or item_id in st.session_state.txt_files_dict)]
    
    # Update file_order to remove invalid references
    if len(valid_file_order) != len(st.session_state.file_order):
        st.session_state.file_order = valid_file_order
    
    # Display files with reordering controls
    for i, item_id in enumerate(valid_file_order):
        col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
        
        with col1:
            st.write(f"**{i+1}.** {item_id}")
        
        with col2:
            if st.button("⬆️ Up", key=f"up_{i}", disabled=(i == 0)):
                if i > 0:
                    valid_file_order[i], valid_file_order[i-1] = valid_file_order[i-1], valid_file_order[i]
                    st.session_state.file_order = valid_file_order
                    st.rerun()
        
        with col3:
            if st.button("⬇️ Down", key=f"down_{i}", disabled=(i == len(valid_file_order) - 1)):
                if i < len(valid_file_order) - 1:
                    valid_file_order[i], valid_file_order[i+1] = valid_file_order[i+1], valid_file_order[i]
                    st.session_state.file_order = valid_file_order
                    st.rerun()
        
        with col4:
            if st.button("🗑️ Remove", key=f"remove_{i}"):
                valid_file_order.remove(item_id)
                st.session_state.file_order = valid_file_order
                if item_id in st.session_state.uploaded_files_dict:
                    del st.session_state.uploaded_files_dict[item_id]
                elif item_id in st.session_state.txt_files_dict:
                    del st.session_state.txt_files_dict[item_id]
                st.rerun()

# Get ordered list of files (both pptx and txt)
ordered_items = []
for item_id in st.session_state.file_order:
    if item_id in st.session_state.uploaded_files_dict:
        # Get file bytes from session state and create BytesIO object
        file_info = st.session_state.uploaded_files_dict[item_id]
        # Create a fresh BytesIO object for each file and ensure it's at the start
        file_bytes_io = BytesIO(file_info['bytes'])
        file_bytes_io.seek(0)  # Ensure we're at the start
        ordered_items.append(('pptx', file_bytes_io))
    elif item_id in st.session_state.txt_files_dict:
        # Get file bytes from session state and create BytesIO object
        file_info = st.session_state.txt_files_dict[item_id]
        file_bytes_io = BytesIO(file_info['bytes'])
        file_bytes_io.seek(0)  # Ensure we're at the start
        ordered_items.append(('txt', file_bytes_io))

# Check if we have any content to merge
has_content = len(ordered_items) > 0

if has_content and st.button("Merge PowerPoints"):
    try:
        merged_presentation = Presentation()
        
        # Set slide dimensions to 16:9 Widescreen aspect ratio
        # Width: 13.33 inches, Height: 7.5 inches
        merged_presentation.slide_width = Inches(13.33)  # 16:9 Widescreen width
        merged_presentation.slide_height = Inches(7.5)  # 16:9 Widescreen height
        
        # Remove default empty slide
        if merged_presentation.slides:
            merged_presentation.slides.remove(merged_presentation.slides[0])
        
        slide_index = 0
        
        for item_type, item_data in ordered_items:
            if item_type == 'pptx':
                # Process PowerPoint file
                # Ensure BytesIO is at the start before creating Presentation
                item_data.seek(0)
                prs = Presentation(item_data)
                first_slide_found = False  # Track if we've found the first slide with text in this file
                first_all_caps_found = False  # Track if we've found the first all-caps slide in this file
                
                for slide in prs.slides:
                    # Extract text from slide
                    text = extract_text_from_slide(slide)
                    
                    if text:  # Only create slide if there's text
                        # Check if text is all caps
                        is_all_caps_text = is_all_caps(text)
                        
                        # Determine if it's a title slide:
                        # - First slide of each PowerPoint file (regardless of caps)
                        # - OR first all-caps slide in each PowerPoint file
                        is_title = False
                        if not first_slide_found:
                            # First slide with text in this file is always a title
                            is_title = True
                            first_slide_found = True
                        elif is_all_caps_text and not first_all_caps_found:
                            # First all-caps slide in this file is also a title
                            is_title = True
                            first_all_caps_found = True
                        
                        # Create formatted slide
                        create_formatted_slide(merged_presentation, text, is_title, st.session_state.title_color, st.session_state.verse_color, st.session_state.title_font_size, st.session_state.verse_font_size, st.session_state.title_font, st.session_state.verse_font, st.session_state.background_image)
                        slide_index += 1
            
            elif item_type == 'txt':
                # Process text file
                item_data.seek(0)  # Ensure we're at the start
                txt_content = item_data.read()
                slides = parse_txt_file(txt_content)
                
                for slide_data in slides:
                    # Create title slide if title exists
                    if slide_data['title']:
                        create_formatted_slide(merged_presentation, slide_data['title'], True, st.session_state.title_color, st.session_state.verse_color, st.session_state.title_font_size, st.session_state.verse_font_size, st.session_state.title_font, st.session_state.verse_font, st.session_state.background_image)
                        slide_index += 1
                    
                    # Create verse slide(s) if verses exist
                    if slide_data['verses']:
                        verses_text = '\n'.join(slide_data['verses'])
                        create_formatted_slide(merged_presentation, verses_text, False, st.session_state.title_color, st.session_state.verse_color, st.session_state.title_font_size, st.session_state.verse_font_size, st.session_state.title_font, st.session_state.verse_font, st.session_state.background_image)
                        slide_index += 1
        
        output = BytesIO()
        merged_presentation.save(output)
        output.seek(0)
        
        st.success("✅ PowerPoints merged successfully!")
        
        st.download_button(
            label="⬇️ Download merged PowerPoint",
            data=output,
            file_name="merged_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"Error merging presentations: {str(e)}")
        st.exception(e)
